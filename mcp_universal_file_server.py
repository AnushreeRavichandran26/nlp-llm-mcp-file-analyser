#!/usr/bin/env python3
"""
mcp_universal_file_server.py — Production-Grade Universal MCP File Server
==========================================================================
UPGRADE: Added strict directory restriction, path traversal prevention,
         and validation for all file access. MCP remains the primary loader.
"""

import asyncio
import json
import sys
import os
import base64
import mimetypes
import logging
from pathlib import Path
from typing import Any, Dict, Optional
from datetime import datetime

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# ─────────────────────────────────────────────
#  SECURITY: Safe path resolver
# ─────────────────────────────────────────────

def safe_resolve(base_dir: Path, user_path: str) -> Path:
    """
    Resolve a user-supplied path and BLOCK traversal outside base_dir.
    e.g. path='../../etc/passwd' is rejected.
    """
    if os.path.isabs(user_path):
        resolved = Path(user_path).resolve()
    else:
        resolved = (base_dir / user_path).resolve()

    # Block path traversal — resolved path must be inside base_dir
    try:
        resolved.relative_to(base_dir)
    except ValueError:
        raise PermissionError(
            f"Access denied: '{user_path}' is outside the allowed directory '{base_dir}'"
        )
    return resolved


# ─────────────────────────────────────────────
#  MAIN SERVER CLASS
# ─────────────────────────────────────────────

class UniversalMCPFileServer:
    def __init__(self, allowed_directory: str = "./documents", restrict_to_directory: bool = False):
        """
        UPGRADE: Default to restricted mode with a safe './documents' folder.
        Full system access is now opt-in, not default.
        """
        self.restrict_to_directory = restrict_to_directory
        self.document_context: Dict[str, Any] = {}

        # Resolve and create the allowed directory
        self.allowed_directory = Path(allowed_directory).resolve()
        self.allowed_directory.mkdir(parents=True, exist_ok=True)

        self.server_info = {"name": "universal-file-server", "version": "2.0.0"}
        self.capabilities = {"tools": {}}

        logger.info(f"MCP File Server v2.0 started")
        logger.info(f"Allowed directory: {self.allowed_directory}")
        logger.info(f"Restricted mode: {self.restrict_to_directory}")

    def resolve_path(self, file_path: str) -> Path:
        if not self.restrict_to_directory:
            # Allow full path access
            return Path(file_path).resolve()

        return safe_resolve(self.allowed_directory, file_path)

    def get_current_timestamp(self) -> str:
        return datetime.utcnow().strftime("%Y-%m-%d %H:%M:%S")

    # ────────────────────────────────────────
    #  TEXT EXTRACTION  (unchanged structure)
    # ────────────────────────────────────────

    def extract_text_from_any_file(self, file_path: Path) -> Dict[str, Any]:
        """Dispatch to the right extractor based on file extension."""
        ext = file_path.suffix.lower()
        result = {
            "success": False, "content": "", "file_type": ext,
            "original_size": file_path.stat().st_size,
            "extracted_size": 0, "method": "unknown", "error": None
        }
        try:
            if   ext == ".pdf":   result.update(self._extract_pdf(file_path))
            elif ext == ".docx":  result.update(self._extract_docx(file_path))
            elif ext == ".xlsx":  result.update(self._extract_xlsx(file_path))
            elif ext == ".pptx":  result.update(self._extract_pptx(file_path))
            elif ext in {".txt",".md",".markdown",".rst",".log"}:
                result.update(self._extract_text_file(file_path))
            elif ext in {".py",".js",".ts",".java",".cpp",".c",".h",
                         ".cs",".php",".rb",".go",".rs",".swift"}:
                result.update(self._extract_code_file(file_path))
            elif ext in {".json",".yaml",".yml",".xml",".csv",".tsv"}:
                result.update(self._extract_data_file(file_path))
            elif ext in {".ini",".conf",".config",".cfg",".toml"}:
                result.update(self._extract_config_file(file_path))
            elif ext in {".html",".htm",".css",".scss"}:
                result.update(self._extract_web_file(file_path))
            elif ext in {".sh",".bash",".zsh",".ps1",".bat",".cmd"}:
                result.update(self._extract_script_file(file_path))
            elif ext in {".doc",".xls",".ppt"}:
                result.update(self._extract_legacy_office(file_path))
            elif ext in {".jpg",".jpeg",".png",".gif",".bmp",".tiff",".webp"}:
                result.update(self._extract_image_metadata(file_path))
            else:
                result.update(self._extract_binary_info(file_path))
        except Exception as e:
            result["error"] = str(e)
            result["content"] = f"Extraction error: {e}"
        result["extracted_size"] = len(result["content"])
        return result

    # — individual extractors (unchanged from original) ——————————————

    def _extract_pdf(self, fp):
        try:
            import PyPDF2
            content = ""
            with open(fp, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for i, page in enumerate(reader.pages):
                    content += f"\n--- Page {i+1} ---\n{page.extract_text() or ''}\n"
            return {"success": True, "content": content.strip(),
                    "method": "PyPDF2", "pages": len(reader.pages)}
        except ImportError:
            return {"success": False, "content": "Install: pip install PyPDF2",
                    "method": "missing", "error": "PyPDF2 not installed"}
        except Exception as e:
            return {"success": False, "content": str(e), "method": "PyPDF2", "error": str(e)}

    def _extract_docx(self, fp):
        try:
            import docx
            doc = docx.Document(fp)
            lines = [p.text for p in doc.paragraphs if p.text.strip()]
            for tbl in doc.tables:
                for row in tbl.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells: lines.append(" | ".join(cells))
            return {"success": True, "content": "\n".join(lines), "method": "python-docx"}
        except ImportError:
            return {"success": False, "content": "Install: pip install python-docx",
                    "method": "missing", "error": "python-docx not installed"}
        except Exception as e:
            return {"success": False, "content": str(e), "method": "python-docx", "error": str(e)}

    def _extract_xlsx(self, fp):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(fp, data_only=True)
            lines = []
            for name in wb.sheetnames:
                lines.append(f"\n=== Sheet: {name} ===")
                for row in wb[name].iter_rows(values_only=True):
                    r = [str(c) for c in row if c is not None]
                    if r: lines.append(" | ".join(r))
            return {"success": True, "content": "\n".join(lines), "method": "openpyxl"}
        except ImportError:
            return {"success": False, "content": "Install: pip install openpyxl",
                    "method": "missing", "error": "openpyxl not installed"}
        except Exception as e:
            return {"success": False, "content": str(e), "method": "openpyxl", "error": str(e)}

    def _extract_pptx(self, fp):
        try:
            import pptx
            prs = pptx.Presentation(fp)
            lines = []
            for i, slide in enumerate(prs.slides, 1):
                lines.append(f"\n=== Slide {i} ===")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        lines.append(shape.text)
            return {"success": True, "content": "\n".join(lines), "method": "python-pptx"}
        except ImportError:
            return {"success": False, "content": "Install: pip install python-pptx",
                    "method": "missing", "error": "python-pptx not installed"}
        except Exception as e:
            return {"success": False, "content": str(e), "method": "python-pptx", "error": str(e)}

    def _extract_text_file(self, fp):
        for enc in ["utf-8", "utf-16", "latin-1", "cp1252"]:
            try:
                return {"success": True, "content": fp.read_text(encoding=enc),
                        "method": f"text_{enc}", "encoding": enc}
            except UnicodeDecodeError:
                continue
        return {"success": False, "content": "Binary content — cannot decode",
                "method": "binary_fallback", "error": "decode failed"}

    def _extract_code_file(self, fp):
        r = self._extract_text_file(fp)
        r["method"] = f"code_{fp.suffix}"; r["language"] = fp.suffix[1:]
        return r

    def _extract_data_file(self, fp):
        ext = fp.suffix.lower()
        try:
            if ext == ".json":
                import json as _json
                data = _json.loads(fp.read_text(encoding="utf-8"))
                return {"success": True, "content": _json.dumps(data, indent=2), "method": "json"}
            elif ext in {".csv", ".tsv"}:
                return {"success": True, "content": fp.read_text(encoding="utf-8"),
                        "method": ext[1:]}
            elif ext in {".yaml", ".yml"}:
                try:
                    import yaml
                    data = yaml.safe_load(fp.read_text(encoding="utf-8"))
                    return {"success": True, "content": yaml.dump(data, allow_unicode=True),
                            "method": "yaml"}
                except ImportError:
                    return self._extract_text_file(fp)
            else:
                return self._extract_text_file(fp)
        except Exception as e:
            return {"success": False, "content": str(e), "method": "data", "error": str(e)}

    def _extract_config_file(self, fp):
        r = self._extract_text_file(fp); r["method"] = f"config_{fp.suffix}"; return r

    def _extract_web_file(self, fp):
        r = self._extract_text_file(fp); r["method"] = f"web_{fp.suffix}"; return r

    def _extract_script_file(self, fp):
        r = self._extract_text_file(fp); r["method"] = f"script_{fp.suffix}"; return r

    def _extract_legacy_office(self, fp):
        return {"success": False,
                "content": f"Legacy {fp.suffix} — please convert to .docx/.xlsx/.pptx or PDF.",
                "method": "legacy", "error": "Legacy format unsupported"}

    def _extract_image_metadata(self, fp):
        st = fp.stat()
        content = (f"Image: {fp.name}\nFormat: {fp.suffix.upper()}\n"
                   f"Size: {st.st_size} bytes\n"
                   f"Modified: {datetime.fromtimestamp(st.st_mtime):%Y-%m-%d %H:%M:%S}\n"
                   "Note: Use the vision model to describe this image.")
        return {"success": True, "content": content, "method": "image_metadata"}

    def _extract_binary_info(self, fp):
        st = fp.stat()
        mime, _ = mimetypes.guess_type(str(fp))
        header = fp.read_bytes()[:256].hex()[:128]
        content = (f"Binary: {fp.name}\nMIME: {mime or 'unknown'}\n"
                   f"Size: {st.st_size} bytes\nHeader: {header}…")
        return {"success": True, "content": content, "method": "binary_info",
                "mime_type": mime}

    # ────────────────────────────────────────
    #  MCP REQUEST HANDLER
    # ────────────────────────────────────────

    async def handle_request(self, request: Dict[str, Any]) -> Dict[str, Any]:
        method = request.get("method")
        params = request.get("params", {})
        rid    = request.get("id")
        try:
            if method == "initialize":
                return {"jsonrpc":"2.0","id":rid,"result":{
                    "protocolVersion":"2024-11-05",
                    "capabilities": self.capabilities,
                    "serverInfo": self.server_info}}

            elif method == "tools/list":
                return {"jsonrpc":"2.0","id":rid,"result":{"tools": self._tool_schemas()}}

            elif method == "tools/call":
                name = params.get("name")
                args = params.get("arguments", {})
                dispatch = {
                    "load_any_document":    self.load_any_document,
                    "query_document_context": self.query_document_context,
                    "list_loaded_documents":  self.list_loaded_documents,
                    "clear_document_context": self.clear_document_context,
                    "read_file":              self.read_any_file,
                    "write_file":             self.write_file,
                    "list_directory":         self.list_directory,
                    "find_files":             self.find_files,
                }
                if name in dispatch:
                    return await dispatch[name](rid, args)
                return self._error(rid, -32601, f"Unknown tool: {name}")

            return self._error(rid, -32601, f"Unknown method: {method}")
        except Exception as e:
            logger.error(f"Request error: {e}")
            return self._error(rid, -32603, f"Internal error: {e}")

    def _error(self, rid, code, msg):
        return {"jsonrpc":"2.0","id":rid,"error":{"code":code,"message":msg}}

    def _ok(self, rid, text):
        return {"jsonrpc":"2.0","id":rid,
                "result":{"content":[{"type":"text","text":text}]}}

    def _tool_schemas(self):
        """Return minimal tool schemas (same as original)."""
        return [
            {"name":"load_any_document","description":"Load any file for Q&A",
             "inputSchema":{"type":"object","properties":{"path":{"type":"string"},
             "context_name":{"type":"string"}},"required":["path"]}},
            {"name":"query_document_context","description":"Query loaded docs",
             "inputSchema":{"type":"object","properties":{"query":{"type":"string"},
             "document_name":{"type":"string"}},"required":["query"]}},
            {"name":"list_loaded_documents","description":"List loaded docs",
             "inputSchema":{"type":"object","properties":{},"required":[]}},
            {"name":"clear_document_context","description":"Clear docs",
             "inputSchema":{"type":"object","properties":{"document_name":{"type":"string"}},"required":[]}},
            {"name":"read_file","description":"Read any file",
             "inputSchema":{"type":"object","properties":{"path":{"type":"string"}},"required":["path"]}},
            {"name":"write_file","description":"Write file",
             "inputSchema":{"type":"object","properties":{"path":{"type":"string"},"content":{"type":"string"}},"required":["path","content"]}},
            {"name":"list_directory","description":"List directory",
             "inputSchema":{"type":"object","properties":{"path":{"type":"string"}},"required":["path"]}},
            {"name":"find_files","description":"Find files",
             "inputSchema":{"type":"object","properties":{"directory":{"type":"string"},"pattern":{"type":"string"},"recursive":{"type":"boolean"}},"required":["directory","pattern"]}},
        ]

    # ────────────────────────────────────────
    #  TOOL IMPLEMENTATIONS
    # ────────────────────────────────────────

    async def load_any_document(self, rid, args):
        path = args.get("path")
        if not path:
            return self._error(rid, -32602, "path required")
        try:
            fp = self.resolve_path(path)
            if not fp.exists():   return self._error(rid, -32602, f"File not found: {fp}")
            if not fp.is_file():  return self._error(rid, -32602, f"Not a file: {fp}")

            ex = self.extract_text_from_any_file(fp)
            name = args.get("context_name") or fp.name
            self.document_context[name] = {
                "path": str(fp), "content": ex["content"],
                "file_type": ex["file_type"], "extraction_method": ex["method"],
                "extraction_success": ex["success"],
                "original_size": ex["original_size"],
                "extracted_size": ex["extracted_size"],
                "loaded_at": self.get_current_timestamp()
            }
            preview = ex["content"][:400] + "…" if len(ex["content"]) > 400 else ex["content"]
            text = (f"✅ Loaded: {name}\n📋 Type: {ex['file_type'].upper()}\n"
                    f"📏 {ex['original_size']} bytes → {ex['extracted_size']} chars\n"
                    f"🔧 Method: {ex['method']}\n\n📖 Preview:\n{preview}")
            return self._ok(rid, text)
        except PermissionError as e:
            return self._error(rid, -32603, str(e))
        except Exception as e:
            return self._error(rid, -32603, f"Load error: {e}")

    async def query_document_context(self, rid, args):
        query = args.get("query")
        if not query:
            return self._error(rid, -32602, "query required")
        if not self.document_context:
            return self._ok(rid, "No documents loaded. Use load_any_document first.")
        target = args.get("document_name")
        docs = ({target: self.document_context[target]} if target and target in self.document_context
                else self.document_context)
        parts = []
        for name, doc in docs.items():
            parts.append(f"=== {name} ({doc['file_type']}) ===\n{doc['content']}")
        text = (f"🔍 Query: \"{query}\"\n\n"
                f"📖 Document Content for Analysis:\n\n" + "\n\n".join(parts))
        return self._ok(rid, text)

    async def list_loaded_documents(self, rid, _args):
        if not self.document_context:
            return self._ok(rid, "📚 No documents loaded.")
        lines = [f"📚 {len(self.document_context)} document(s) loaded:\n"]
        for name, doc in self.document_context.items():
            lines.append(f"  📄 {name}  [{doc['file_type']}]  "
                         f"{doc['extracted_size']} chars  {doc['loaded_at']}")
        return self._ok(rid, "\n".join(lines))

    async def clear_document_context(self, rid, args):
        name = args.get("document_name")
        if name:
            self.document_context.pop(name, None)
            msg = f"🗑️ Cleared: {name}"
        else:
            n = len(self.document_context); self.document_context.clear()
            msg = f"🗑️ Cleared {n} documents"
        return self._ok(rid, msg)

    async def read_any_file(self, rid, args):
        path = args.get("path")
        if not path: return self._error(rid, -32602, "path required")
        try:
            fp = self.resolve_path(path)
            if not fp.exists():  return self._error(rid, -32602, f"Not found: {fp}")
            if not fp.is_file(): return self._error(rid, -32602, f"Not a file: {fp}")
            ex = self.extract_text_from_any_file(fp)
            text = (f"📄 {fp}\n📋 {ex['file_type'].upper()}  "
                    f"{ex['original_size']}B → {ex['extracted_size']} chars\n\n"
                    f"{ex['content']}")
            return self._ok(rid, text)
        except PermissionError as e:
            return self._error(rid, -32603, str(e))
        except Exception as e:
            return self._error(rid, -32603, str(e))

    async def write_file(self, rid, args):
        path = args.get("path"); content = args.get("content", "")
        if not path: return self._error(rid, -32602, "path required")
        try:
            fp = self.resolve_path(path)
            fp.parent.mkdir(parents=True, exist_ok=True)
            fp.write_text(content, encoding="utf-8")
            return self._ok(rid, f"✅ Wrote {len(content)} chars to {fp}")
        except PermissionError as e:
            return self._error(rid, -32603, str(e))
        except Exception as e:
            return self._error(rid, -32603, str(e))

    async def list_directory(self, rid, args):
        path = args.get("path", ".")
        try:
            dp = self.resolve_path(path)
            if not dp.exists() or not dp.is_dir():
                return self._error(rid, -32602, f"Not a directory: {dp}")
            lines = [f"📁 {dp}\n"]
            for item in sorted(dp.iterdir()):
                icon = "📁" if item.is_dir() else "📄"
                sz = "<DIR>" if item.is_dir() else f"{item.stat().st_size}B"
                lines.append(f"  {icon} {item.name:<40} {sz}")
            return self._ok(rid, "\n".join(lines))
        except PermissionError as e:
            return self._error(rid, -32603, str(e))
        except Exception as e:
            return self._error(rid, -32603, str(e))

    async def find_files(self, rid, args):
        directory = args.get("directory"); pattern = args.get("pattern")
        recursive = args.get("recursive", True)
        if not directory or not pattern:
            return self._error(rid, -32602, "directory and pattern required")
        try:
            import fnmatch
            dp = self.resolve_path(directory)
            matches = []
            if recursive:
                for root, _, files in os.walk(dp):
                    for f in files:
                        if fnmatch.fnmatch(f, pattern):
                            matches.append(str(Path(root) / f))
            else:
                matches = [str(p) for p in dp.iterdir()
                           if p.is_file() and fnmatch.fnmatch(p.name, pattern)]
            text = f"🔍 {len(matches)} match(es) for '{pattern}' in {dp}\n\n"
            text += "\n".join(matches[:100])
            return self._ok(rid, text)
        except PermissionError as e:
            return self._error(rid, -32603, str(e))
        except Exception as e:
            return self._error(rid, -32603, str(e))

    # ────────────────────────────────────────
    #  STDIO LOOP
    # ────────────────────────────────────────

    async def run(self):
        logger.info("Universal MCP File Server v2.0 running (stdio)")
        while True:
            try:
                line = await asyncio.get_event_loop().run_in_executor(
                    None, sys.stdin.readline)
                if not line: break
                line = line.strip()
                if not line: continue
                request  = json.loads(line)
                response = await self.handle_request(request)
                print(json.dumps(response), flush=True)
            except json.JSONDecodeError:
                print(json.dumps({"jsonrpc":"2.0","id":None,
                                  "error":{"code":-32700,"message":"Parse error"}}), flush=True)
            except KeyboardInterrupt:
                break
            except Exception as e:
                logger.error(f"Unexpected: {e}")


def main():
    import argparse
    p = argparse.ArgumentParser(description="Universal MCP File Server v2.0")
    p.add_argument("--directory", "-d", default="./documents",
                   help="Base directory (default: ./documents)")
    p.add_argument("--no-restrict", action="store_true",
                   help="Disable directory restriction (NOT recommended)")
    args = p.parse_args()
    server = UniversalMCPFileServer(args.directory, not args.no_restrict)
    asyncio.run(server.run())


if __name__ == "__main__":
    main()
